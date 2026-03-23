"""
Builds a .pptx with LARGER text (16-20pt body) and simplified language.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


CHARTS = "charts"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD = RGBColor(0x1E, 0x29, 0x3B)
ACCENT  = RGBColor(0x81, 0x8C, 0xF8)
CYAN    = RGBColor(0x22, 0xD3, 0xEE)
WHITE   = RGBColor(0xE2, 0xE8, 0xF0)
MUTED   = RGBColor(0x94, 0xA3, 0xB8)
DIM     = RGBColor(0x64, 0x74, 0x8B)
BORDER  = RGBColor(0x33, 0x41, 0x55)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def tb(slide, left, top, width, height, text, size=18,
       color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    return txBox


def rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def img(slide, path, left, top, width=None, height=None):
    if width and height:
        slide.shapes.add_picture(path, left, top, width, height)
    elif width:
        slide.shapes.add_picture(path, left, top, width=width)
    else:
        slide.shapes.add_picture(path, left, top)


def stat(slide, left, top, number, label):
    rect(slide, left, top, Inches(2.3), Inches(1.2),
         RGBColor(0x16, 0x1D, 0x35), ACCENT)
    tb(slide, left, top + Inches(0.1), Inches(2.3), Inches(0.6),
       str(number), 30, ACCENT, True, PP_ALIGN.CENTER)
    tb(slide, left, top + Inches(0.7), Inches(2.3), Inches(0.4),
       label, 14, MUTED, False, PP_ALIGN.CENTER)


def card(slide, left, top, width, height, title, body):
    rect(slide, left, top, width, height, BG_CARD, BORDER)
    tb(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
       title, 18, ACCENT, True)
    tb(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
       body, 16, MUTED)


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# ===== SLIDE 1: TITLE =====
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
tb(s, Inches(1), Inches(2.0), Inches(11.3), Inches(1.5),
   "Google Sustainability Analytics\nSuggestion System",
   44, ACCENT, True, PP_ALIGN.CENTER)
tb(s, Inches(1.5), Inches(3.9), Inches(10.3), Inches(0.7),
   "Using Machine Learning to Help Buildings Save Energy",
   24, MUTED, False, PP_ALIGN.CENTER)
tb(s, Inches(1.5), Inches(4.8), Inches(10.3), Inches(0.5),
   "AI/ML Project Presentation", 18, DIM, False, PP_ALIGN.CENTER)


# ===== SLIDE 2: PROBLEM =====
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
tb(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
   "What Problem Are We Solving?", 34, WHITE, True)

pts = [
    ("The Problem:", "Buildings waste a huge amount of electricity, but owners have no idea WHERE the waste is happening or WHAT to do about it."),
    ("Why Its Hard:", "Every building is different. An office wastes energy differently from a house. Generic tips like 'turn off lights' dont work."),
    ("Our Solution:", "Automatically group similar buildings together using ML, then give each group specific energy-saving advice."),
]
for i, (label, body) in enumerate(pts):
    y = Inches(1.3) + Inches(1.3) * i
    rect(s, Inches(0.5), y, Inches(12.3), Inches(1.1), BG_CARD, BORDER)
    tb(s, Inches(0.8), y + Inches(0.1), Inches(11.7), Inches(0.9),
       f"{label} {body}", 18, MUTED)

stat(s, Inches(1.5), Inches(5.6), "30%+", "Energy by Buildings")
stat(s, Inches(5.5), Inches(5.6), "55%", "Electricity to Buildings")
stat(s, Inches(9.5), Inches(5.6), "15-30%", "Savings Possible")


# ===== SLIDE 3: HOW IT WORKS =====
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
tb(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
   "How Does It Work? (6 Steps)", 34, WHITE, True)

steps = [
    ("Step 1: Get Data", "Created realistic fake electricity data for 25 buildings. Smart meter readings every 15 minutes for 30 days."),
    ("Step 2: Describe Buildings", "From raw readings, calculated 11 numbers per building (average usage, peak, behavior patterns, etc)."),
    ("Step 3: Normalize", "Scaled all numbers to same range so no single feature dominates the analysis."),
    ("Step 4: Compress (PCA)", "Compressed 11 numbers into 4 using PCA, keeping 95% of useful information."),
    ("Step 5: Group (K-Means)", "Algorithm sorted 25 buildings into 8 groups based on energy pattern similarity."),
    ("Step 6: Recommend", "Rules check each group's profile and suggest specific energy saving actions."),
]
cw = Inches(3.9)
ch = Inches(2.5)
for i, (title, body) in enumerate(steps):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + col * (cw + Inches(0.3))
    y = Inches(1.2) + row * (ch + Inches(0.2))
    card(s, x, y, cw, ch, title, body)


# ===== SLIDE 4: FEATURES =====
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
tb(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
   "The 11 Numbers Describing Each Building", 34, WHITE, True)

features = [
    ("avg_kw", "Average electricity usage"),
    ("peak_kw", "Highest spike in usage"),
    ("base_kw", "Always-on background usage"),
    ("load_factor", "Steady or spiky? (1=steady)"),
    ("variability_cv", "How unpredictable is usage?"),
    ("peak_hour_mode", "Time of day for peak usage"),
    ("peak_to_offpeak", "Day vs night usage ratio"),
    ("wkend_wkday", "Weekend vs weekday ratio"),
    ("hvac_share", "Heating/cooling percentage"),
    ("standby_share", "Always-on device percentage"),
    ("shiftable_share", "Movable to cheaper hours"),
]
for i, (feat, desc) in enumerate(features):
    y = Inches(1.2) + i * Inches(0.52)
    tb(s, Inches(0.5), y, Inches(2.5), Inches(0.45), feat, 16, CYAN, True)
    tb(s, Inches(3.1), y, Inches(3.5), Inches(0.45), desc, 16, MUTED)

img(s, os.path.join(CHARTS, "feature_importance.png"),
    Inches(7.0), Inches(1.0), width=Inches(6.0))
tb(s, Inches(7.0), Inches(6.5), Inches(6.0), Inches(0.5),
   "Which features matter most for separating buildings", 14, DIM, False, PP_ALIGN.CENTER)


# ===== SLIDE 5: PCA =====
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
tb(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
   "PCA - Simplifying The Data", 34, WHITE, True)

pca = [
    ("What is PCA?", "Imagine describing someone with 11 facts. PCA finds that just 4 combined facts capture 95% of the same info. It compresses data."),
    ("Why use it?", "Grouping with 11 numbers is noisy. 4 numbers give cleaner groups. Also lets us draw things on a chart."),
    ("PC1 = Building Size", "Captures 57.8% of info. Big consumers on one side, small ones on the other."),
    ("PC2 = Behavior", "Captures 21.2%. Separates by patterns like weekday vs weekend usage."),
]
for i, (label, body) in enumerate(pca):
    y = Inches(1.2) + i * Inches(1.3)
    rect(s, Inches(0.5), y, Inches(6.0), Inches(1.15), BG_CARD, BORDER)
    tb(s, Inches(0.75), y + Inches(0.08), Inches(5.5), Inches(1.0),
       f"{label}  {body}", 16, MUTED)

img(s, os.path.join(CHARTS, "pca_variance.png"),
    Inches(6.8), Inches(1.0), width=Inches(6.0))
tb(s, Inches(6.8), Inches(6.3), Inches(6.0), Inches(0.5),
   "4 bars cross the 95% line = 4 components are enough", 15, DIM, False, PP_ALIGN.CENTER)


# ===== SLIDE 6: PCA LOADINGS =====
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
tb(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
   "Which Features Drive Each Component?", 34, WHITE, True)
img(s, os.path.join(CHARTS, "pca_loadings.png"),
    Inches(0.5), Inches(1.1), width=Inches(12.3))
rect(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.5), BG_CARD, BORDER)
tb(s, Inches(1.0), Inches(5.6), Inches(11.3), Inches(1.3),
   "Taller bars = that feature influences that component more. "
   "PC1 (purple) is driven by avg_kw and standby_share = building size. "
   "PC2 (teal) is driven by weekend ratio and variability = behavior pattern.",
   18, MUTED)


# ===== SLIDE 7: K-MEANS =====
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
tb(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
   "K-Means - Grouping Similar Buildings", 34, WHITE, True)

km = [
    ("What is it?", "An algorithm that sorts buildings into groups. Similar buildings end up together. No manual labeling needed."),
    ("How it works:", "Place 8 random centers. Each building joins nearest center. Move centers to group middle. Repeat until stable."),
    ("Why 8 groups?", "Tried 2 to 10. Measured quality for each. k=8 scored highest quality (0.598 out of 1.0)."),
    ("Quality check:", "Silhouette Score asks each building 'are you closer to your group or the wrong group?' 0.598 = good."),
]
for i, (label, body) in enumerate(km):
    y = Inches(1.2) + i * Inches(1.2)
    rect(s, Inches(0.5), y, Inches(6.0), Inches(1.05), BG_CARD, BORDER)
    tb(s, Inches(0.75), y + Inches(0.08), Inches(5.5), Inches(0.9),
       f"{label}  {body}", 16, MUTED)

img(s, os.path.join(CHARTS, "k_sweep.png"),
    Inches(6.8), Inches(1.0), width=Inches(6.0))
tb(s, Inches(6.8), Inches(6.3), Inches(6.0), Inches(0.5),
   "Red drops = tighter clusters. Teal peaks at k=8 = best quality.", 15, DIM, False, PP_ALIGN.CENTER)


# ===== SLIDE 8: CLUSTER VIZ =====
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
tb(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
   "The Result - 8 Building Groups", 34, WHITE, True)

cv = [
    ("What is this chart?", "Each dot = one building. X = size, Y = behavior. Colors = which group it belongs to."),
    ("What does it show?", "Same-colored dots cluster together. This means the algorithm found real patterns, not random noise."),
    ("Why it matters:", "This proves our approach worked. Buildings in each color genuinely share energy patterns."),
]
for i, (label, body) in enumerate(cv):
    y = Inches(1.2) + i * Inches(1.5)
    rect(s, Inches(0.5), y, Inches(5.8), Inches(1.3), BG_CARD, BORDER)
    tb(s, Inches(0.75), y + Inches(0.08), Inches(5.3), Inches(1.1),
       f"{label}  {body}", 17, MUTED)

img(s, os.path.join(CHARTS, "pca_clusters.png"),
    Inches(6.6), Inches(1.0), width=Inches(6.3))


# ===== SLIDE 9: EVALUATION =====
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
tb(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
   "How Good Are Our Groups?", 34, WHITE, True)

ev = [
    ("Overall: 0.598", "Scale is -1 to +1. Above 0.5 = well-separated groups. Ours is solidly good."),
    ("Best groups:", "Clusters 3, 4, 5, 7 scored above 0.7. These buildings are very clearly grouped."),
    ("Weaker ones:", "Clusters 0, 1 scored ~0.15. Small buildings that look somewhat similar to each other."),
]
for i, (label, body) in enumerate(ev):
    y = Inches(1.2) + i * Inches(1.2)
    rect(s, Inches(0.5), y, Inches(5.8), Inches(1.05), BG_CARD, BORDER)
    tb(s, Inches(0.75), y + Inches(0.08), Inches(5.3), Inches(0.9),
       f"{label}  {body}", 17, MUTED)

stat(s, Inches(0.6), Inches(5.5), "0.598", "Quality Score")
stat(s, Inches(3.3), Inches(5.5), "8", "Groups Found")
stat(s, Inches(6.0), Inches(5.5), "94.8%", "Info Kept by PCA")

img(s, os.path.join(CHARTS, "silhouette_box.png"),
    Inches(6.6), Inches(1.0), width=Inches(6.3))


# ===== SLIDE 10: HEATMAP =====
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
tb(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
   "What Makes Each Group Different?", 34, WHITE, True)
img(s, os.path.join(CHARTS, "cluster_heatmap.png"),
    Inches(0.3), Inches(1.1), width=Inches(12.7))
tb(s, Inches(0.3), Inches(5.1), Inches(12.7), Inches(0.5),
   "Red = high, Blue = low. Each row is a group, each column is a feature.", 16, DIM, False, PP_ALIGN.CENTER)

profiles = [
    ("Group 0: Small + Always-On", "Low energy but high standby. Lots of devices running 24/7."),
    ("Group 3: Big HVAC Buildings", "Highest consumption. Dominated by heating and cooling."),
    ("Group 5: Extreme Peaks", "Biggest day-vs-night gap. Needs load shifting the most."),
]
for i, (label, body) in enumerate(profiles):
    x = Inches(0.3) + i * Inches(4.3)
    rect(s, x, Inches(5.7), Inches(4.1), Inches(1.4), BG_CARD, BORDER)
    tb(s, x + Inches(0.15), Inches(5.8), Inches(3.8), Inches(0.35), label, 15, ACCENT, True)
    tb(s, x + Inches(0.15), Inches(6.2), Inches(3.8), Inches(0.7), body, 15, MUTED)


# ===== SLIDE 11: RECOMMENDATIONS =====
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
tb(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
   "What Advice Does The System Give?", 34, WHITE, True)

tb(s, Inches(0.5), Inches(1.1), Inches(7), Inches(0.5),
   "It checks each building's numbers and suggests fixes:", 18, MUTED)

recs = [
    ("Day usage >> night usage", "Move EV charging, dishwasher, laundry to night hours"),
    ("Too many always-on devices", "Use smart power strips to cut standby drain"),
    ("Heating/cooling over 40%", "Adjust thermostat, improve insulation, service HVAC"),
    ("Very unpredictable usage", "Spread out high-power appliances, dont run all at once"),
    ("Peak hits in the evening", "Pre-cool/heat earlier so it doesnt all spike at 7pm"),
]
for i, (cond, rec) in enumerate(recs):
    y = Inches(1.7) + i * Inches(0.95)
    rect(s, Inches(0.5), y, Inches(7.0), Inches(0.85), BG_CARD, BORDER)
    tb(s, Inches(0.7), y + Inches(0.05), Inches(2.9), Inches(0.75), cond, 15, CYAN, True)
    tb(s, Inches(3.7), y + Inches(0.05), Inches(3.5), Inches(0.75), rec, 15, MUTED)

img(s, os.path.join(CHARTS, "cluster_counts.png"),
    Inches(7.8), Inches(1.0), width=Inches(5.0))
tb(s, Inches(7.8), Inches(5.8), Inches(5.0), Inches(0.5),
   "How many buildings ended up in each group", 14, DIM, False, PP_ALIGN.CENTER)


# ===== SLIDE 12: TECH STACK =====
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
tb(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
   "Tools We Used", 34, WHITE, True)

techs = [
    ("Python", "Everything is written in Python"),
    ("pandas", "Data tables, like Excel in code"),
    ("NumPy", "Fast math on arrays of numbers"),
    ("scikit-learn", "ML library: PCA, K-Means, scoring"),
    ("joblib", "Save/load trained models"),
    ("Streamlit", "Python to web dashboard, zero HTML"),
    ("Plotly", "Interactive zoomable charts"),
]
for i, (tech, purpose) in enumerate(techs):
    y = Inches(1.2) + i * Inches(0.65)
    tb(s, Inches(0.5), y, Inches(2.5), Inches(0.55), tech, 20, CYAN, True)
    tb(s, Inches(3.2), y, Inches(4), Inches(0.55), purpose, 18, MUTED)

card(s, Inches(7.2), Inches(1.2), Inches(5.5), Inches(2.5),
     "Live Dashboard (4 Tabs)",
     "1. Summary: all clusters at a glance\n"
     "2. PCA: variance + scatter plots\n"
     "3. K-Means: how we picked k=8\n"
     "4. Deep Dive: pick a building, see tips")
card(s, Inches(7.2), Inches(4.0), Inches(5.5), Inches(1.5),
     "Accepts New Data",
     "Upload a CSV and it classifies new buildings into existing groups.")


# ===== SLIDE 13: CONCLUSION =====
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
tb(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
   "What We Achieved", 34, WHITE, True)

conc = [
    ("End-to-End Pipeline", "Built everything: data generation, feature extraction, PCA, K-Means, web dashboard. All works together."),
    ("Why Unsupervised?", "Real buildings dont have labels. Unsupervised ML discovers patterns on its own. Practical for real-world use."),
    ("Key Numbers", "8 groups found. 0.598 quality score (good). PCA kept 94.8% of info in just 4 numbers."),
    ("Real Impact", "Targeted advice per group can save 15-30% more energy than generic tips."),
]
cw2 = Inches(6.0)
ch2 = Inches(1.8)
for i, (title, body) in enumerate(conc):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + col * (cw2 + Inches(0.3))
    y = Inches(1.2) + row * (ch2 + Inches(0.25))
    card(s, x, y, cw2, ch2, title, body)

stat(s, Inches(0.3), Inches(5.7), "11", "Features")
stat(s, Inches(2.9), Inches(5.7), "8", "Groups")
stat(s, Inches(5.5), Inches(5.7), "0.598", "Quality")
stat(s, Inches(8.1), Inches(5.7), "94.8%", "Info Kept")
stat(s, Inches(10.7), Inches(5.7), "25", "Buildings")


# ===== SLIDE 14: THANK YOU =====
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
tb(s, Inches(1), Inches(2.3), Inches(11.3), Inches(1.2),
   "Thank You!", 52, ACCENT, True, PP_ALIGN.CENTER)
tb(s, Inches(1.5), Inches(3.8), Inches(10.3), Inches(0.7),
   "Google Sustainability Analytics Suggestion System",
   24, MUTED, False, PP_ALIGN.CENTER)
tb(s, Inches(1.5), Inches(4.7), Inches(10.3), Inches(0.5),
   "Questions?", 22, DIM, False, PP_ALIGN.CENTER)


prs.save("presentation.pptx")
print(f"Saved presentation.pptx  |  {len(prs.slides)} slides")
